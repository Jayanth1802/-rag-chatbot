import os
import tempfile
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from docx2txt import process as docx_process
from pptx import Presentation
from langchain.schema import Document

FAISS_PATH = "faiss_index"

def load_pdf(file_path):
    loader = PyPDFLoader(file_path)
    return loader.load()

def load_docx(file_path):
    text = docx_process(file_path)
    return [Document(page_content=text, metadata={"source": file_path})]

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return [Document(page_content=text, metadata={"source": file_path})]

def ingest_uploaded_files(uploaded_files):
    all_documents = []

    for uploaded_file in uploaded_files:
        suffix = os.path.splitext(uploaded_file.name)[1].lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uploaded_file.read())
            tmp_path = tmp.name

        if suffix == ".pdf":
            docs = load_pdf(tmp_path)
        elif suffix == ".docx":
            docs = load_docx(tmp_path)
        elif suffix == ".pptx":
            docs = load_pptx(tmp_path)
        else:
            continue

        all_documents.extend(docs)
        os.unlink(tmp_path)

    if not all_documents:
        return False

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    chunks = splitter.split_documents(all_documents)

    embeddings = HuggingFaceEmbeddings(
        model_name="all-MiniLM-L6-v2"
    )
    vectorstore = FAISS.from_documents(chunks, embeddings)
    vectorstore.save_local(FAISS_PATH)
    return True